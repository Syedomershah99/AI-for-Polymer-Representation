#!/usr/bin/env python3
"""Update PowerPoint presentation with new polymer analysis plots."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Paths
PPT_PATH = "/Users/Omer/Desktop/Research/Polymer Representation.pptx"
PLOTS_DIR = "/Users/Omer/Desktop/Research/plots"
OUTPUT_PATH = "/Users/Omer/Desktop/Research/Polymer Representation_Updated.pptx"

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_image_slide(prs, title, image_path, observations=None):
    """Add a slide with an image and optional observations."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Image
    if os.path.exists(image_path):
        if observations:
            # Image on left, observations on right
            slide.shapes.add_picture(image_path, Inches(0.2), Inches(0.9), width=Inches(6.2))

            # Observations box
            obs_box = slide.shapes.add_textbox(Inches(6.5), Inches(0.9), Inches(3.3), Inches(6))
            tf = obs_box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = "Observations:"
            p.font.size = Pt(16)
            p.font.bold = True

            for obs in observations:
                p = tf.add_paragraph()
                p.text = f"• {obs}"
                p.font.size = Pt(12)
                p.space_before = Pt(6)
        else:
            # Full width image
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(0.9), width=Inches(9))

    return slide

def add_two_image_slide(prs, title, image_path1, image_path2, observations=None):
    """Add a slide with two images side by side."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Images side by side
    if os.path.exists(image_path1):
        slide.shapes.add_picture(image_path1, Inches(0.2), Inches(0.9), width=Inches(4.7))
    if os.path.exists(image_path2):
        slide.shapes.add_picture(image_path2, Inches(5.1), Inches(0.9), width=Inches(4.7))

    # Observations at bottom
    if observations:
        obs_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(9.4), Inches(2))
        tf = obs_box.text_frame
        tf.word_wrap = True
        for i, obs in enumerate(observations):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {obs}"
            p.font.size = Pt(11)

    return slide

def add_text_slide(prs, title, content_items):
    """Add a text-only slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.space_before = Pt(12)

    return slide

def main():
    # Load existing presentation
    prs = Presentation(PPT_PATH)

    # Add section divider
    add_title_slide(prs, "Week 3 Update", "January 6, 2026 - Polymer Representation Analysis")

    # 1. K-Means Clustering Analysis
    add_image_slide(
        prs,
        "K-Means Clustering: Elbow & Silhouette Analysis",
        os.path.join(PLOTS_DIR, "01_kmeans_analysis.png"),
        observations=[
            "Elbow curve shows diminishing returns after k=8-10 clusters",
            "Silhouette score peaks around k=6-8, indicating natural groupings",
            "Optimal k chosen based on balance of cluster compactness and separation",
            "Morgan FP representation captures structural diversity effectively"
        ]
    )

    # 2. Hierarchical Clustering Analysis
    add_two_image_slide(
        prs,
        "Hierarchical Clustering Analysis",
        os.path.join(PLOTS_DIR, "02_hierarchical_elbow_analysis.png"),
        os.path.join(PLOTS_DIR, "02_dendrogram_morgan.png"),
        observations=[
            "Hierarchical elbow analysis confirms k=8-10 as optimal cluster count",
            "Dendrogram reveals clear hierarchical structure in polymer space",
            "Ward linkage provides balanced cluster sizes"
        ]
    )

    # 3. Butina Clustering Cutoff Analysis
    add_image_slide(
        prs,
        "Butina Clustering: Tanimoto Distance Cutoff Analysis",
        os.path.join(PLOTS_DIR, "02_butina_cutoff_analysis.png"),
        observations=[
            "Cutoff of 0.3-0.4 provides good balance of cluster count and coverage",
            "Lower cutoffs create many singleton clusters",
            "Higher cutoffs merge structurally distinct polymers",
            "Silhouette score guides optimal cutoff selection"
        ]
    )

    # 4. Clustering Validation
    add_image_slide(
        prs,
        "Clustering Validation: ARI & NMI Metrics",
        os.path.join(PLOTS_DIR, "03_clustering_validation.png"),
        observations=[
            "ARI measures agreement with polymer class labels",
            "NMI quantifies mutual information between clusters and classes",
            "Higher values indicate clustering aligns with chemical classes",
            "Morgan FP-based clustering shows strong class correspondence"
        ]
    )

    # 5. Supervised Model Performance
    add_image_slide(
        prs,
        "Supervised Learning: 5-Fold CV Accuracy",
        os.path.join(PLOTS_DIR, "04_supervised_accuracy.png"),
        observations=[
            "Random Forest achieves best performance on Morgan FP",
            "Stratified 5-fold CV ensures balanced class distribution",
            "Different representations show varying classification power",
            "Transformer embeddings capture semantic similarity"
        ]
    )

    # 6. UMAP Visualizations
    add_image_slide(
        prs,
        "UMAP: Polymer Space by Class",
        os.path.join(PLOTS_DIR, "05_umap_by_class.png"),
        observations=[
            "Clear separation between polymer classes in UMAP space",
            "Similar chemical families cluster together",
            "Some overlap indicates structural similarities across classes",
            "UMAP preserves local and global structure"
        ]
    )

    add_image_slide(
        prs,
        "UMAP: K-Means Cluster Distribution",
        os.path.join(PLOTS_DIR, "06_umap_by_kmeans.png"),
        observations=[
            "K-means clusters align well with UMAP embedding structure",
            "Cluster boundaries correspond to density variations",
            "Validates clustering quality on structural representations"
        ]
    )

    # 7. Property-colored UMAP
    add_two_image_slide(
        prs,
        "UMAP: Property Distributions (Physical)",
        os.path.join(PLOTS_DIR, "07_umap_by_density.png"),
        os.path.join(PLOTS_DIR, "07_umap_by_bulk_modulus.png"),
        observations=[
            "Density shows gradient patterns across polymer space",
            "Bulk modulus correlates with specific structural regions",
            "Property trends visible in embedding space"
        ]
    )

    add_two_image_slide(
        prs,
        "UMAP: Property Distributions (Thermal & Electrical)",
        os.path.join(PLOTS_DIR, "07_umap_by_thermal_conductivity.png"),
        os.path.join(PLOTS_DIR, "07_umap_by_static_dielectric_const.png"),
        observations=[
            "Thermal conductivity varies with polymer backbone structure",
            "Dielectric constant shows distinct regions in embedding",
            "Structure-property relationships visible in 2D projection"
        ]
    )

    # 8. MACCS Keys Analysis
    add_two_image_slide(
        prs,
        "MACCS Keys: Substructure Distribution",
        os.path.join(PLOTS_DIR, "08_maccs_distribution.png"),
        os.path.join(PLOTS_DIR, "09_maccs_class_heatmap.png"),
        observations=[
            "166 MACCS keys capture standard substructure patterns",
            "Aromatic rings, heteroatoms, chain features most prevalent",
            "Heatmap shows class-specific motif signatures",
            "Enables interpretable motif-based polymer characterization"
        ]
    )

    # 9. Representative Structures
    add_two_image_slide(
        prs,
        "Representative Polymer Structures by Cluster",
        os.path.join(PLOTS_DIR, "08_kmeans_representatives.png"),
        os.path.join(PLOTS_DIR, "09_hierarchical_representatives.png"),
        observations=[
            "Representatives selected as closest to cluster centroid",
            "Captures structural diversity within dataset",
            "K-means and hierarchical methods identify similar representatives",
            "Useful for designing representative experimental sets"
        ]
    )

    # 10. Key Findings Summary
    add_text_slide(
        prs,
        "Key Findings - Week 3",
        [
            "Optimal cluster count: k=8-10 based on elbow/silhouette analysis",
            "Morgan FP (2048-bit) provides strong structural representation",
            "MACCS keys enable interpretable motif-based characterization",
            "Clustering aligns well with polymer class labels (validated via ARI/NMI)",
            "UMAP reveals clear structure-property relationships",
            "Butina clustering with 0.3-0.4 Tanimoto cutoff gives balanced results",
            "Supervised models achieve good classification accuracy with Random Forest",
            "Transformer embeddings (polyBERT) capture semantic similarity"
        ]
    )

    # 11. Next Steps
    add_text_slide(
        prs,
        "Next Steps",
        [
            "Integrate clustering results into Bayesian Optimization pipeline",
            "Use representative polymers for initial experimental screening",
            "Explore active learning with cluster-based acquisition functions",
            "Extend analysis to multi-property optimization",
            "Validate structure-property correlations with DFT calculations",
            "Build property prediction models using combined representations"
        ]
    )

    # Save updated presentation
    prs.save(OUTPUT_PATH)
    print(f"Updated presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
